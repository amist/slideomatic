#-------------------------------------------------------------------------------
# Name:        module1
# Purpose:
#
# Author:      yfodor
#
# Created:     17/06/2015
# Copyright:   (c) yfodor 2015
# Licence:     <your licence>
#-------------------------------------------------------------------------------

#*************************************************************************
# GREAT LINK:
# https://python-pptx.readthedocs.org/en/latest/user/quickstart.html
#*************************************************************************

from pptx import Presentation
from pptx.util import Inches, Pt
import math


def add_title_slide(prs, title_dict):
    '''
    The function get dictionary with the following keys and create title slide:
        title - the title of the presentation
        author - the name of the author - will present in the bottom of the slide
        img_title - path to the image. if picture is not needed this key is not valid
    '''
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_dict["title"]
    subtitle.text = title_dict["author"]
    if(title_dict.has_key("image")):
        top = Inches(1)
        left = Inches(3)
        height = Inches(3)

        pic = slide.shapes.add_picture(title_dict["image"], left, top, height=height)


def add_contact_slide(prs, contact_dict):
    '''
    Optional: if we want to have contact slide
    Yoav suggest to pass it as a regular text if needed
    '''
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Contact Details'

    tf = body_shape.text_frame
    tf.text = contact_dict["name"]
    if contact_dict.has_key("email"):
        p = tf.add_paragraph()
        p.text = contact_dict["email"]
    if contact_dict.has_key("phone"):
        p = tf.add_paragraph()
        p.text = contact_dict["phone"]

def add_text_in_slide(slide, slide_dict):
    '''
    this function get as input the presentation and slide_dict that includes the following relevant key:
    text - the text to be added
    '''
    #blank_slide_layout = prs.slide_layouts[6]
    #slide = prs.slides.add_slide(blank_slide_layout)
    #position of the text-box
    number_of_lines = slide_dict["text"].count('\n') + 2
    text_lines = slide_dict["text"].split('\n')
    max_len = 0
    for line in text_lines:
        if len(line) > max_len:
            max_len = len(line)

    left = top = Inches(1)
    width = Inches(math.ceil(max_len / 6.0))          #6 characters in size 32 for an inche
    height = Inches(math.ceil(number_of_lines / 2)) #2 lines for an inche
    print number_of_lines, width, height
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    #tf.text = slide_dict["text"]
    p = tf.add_paragraph()
    p.text = slide_dict["text"]
    p.font.size = Pt(32)


def parsing_and_create_slides(prs, data_list):
    '''
    this function get the presentation and list of:
    {'text': "some_text", 'img_list':'img_list}
    '''
    for slide_dict in data_list:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        add_text_in_slide(slide, slide_dict)

def make_presentation(title_dict, data_list):
    #title_dict = {"title":"Go Go Slide-O-Matic", "author":"Slide-O-Matic Team", }#"image":"title_img.jpg"}

    #contact_dict = {"name": "The marvelous Team\nTest",  "phone": "10-9", "email": "marvelous@marvel.com"}

    prs = Presentation()

    add_title_slide(prs, title_dict)
    parsing_and_create_slides(prs, data_list)
    #add_text_slide(prs, slide1)
    #add_contact_slide(prs, contact_dict)

    prs.save('test.pptx')

if __name__ == '__main__':
    title_dict = {"title":"Go Go Slide-O-Matic", "author":"Slide-O-Matic Team", }#"image":"title_img.jpg"}
    slide1 = {"text": "* yair fodor\n* is one of the best\n* fuzball legs up player ever"}
    slide2 = {"text": "* alon shaltiel - go to sleep please..."}
    make_presentation(title_dict, [slide1, slide2])

